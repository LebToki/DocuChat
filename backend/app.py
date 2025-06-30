import os
import json
import numpy as np
import torch
from flask import Flask, request, jsonify, session
from flask_cors import CORS
from pdfminer.high_level import extract_text as extract_text_from_pdf
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from transformers import AutoModel, AutoTokenizer, BertForSequenceClassification, TrainingArguments, Trainer
import faiss
from fine_tune_model import ModelFineTuner
from langdetect import detect

# Set environment variable to handle OpenMP error
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"

app = Flask(__name__)
app.secret_key = os.environ.get('SECRET_KEY', 'default_secret_key')
CORS(app, resources={r"/*": {"origins": "*"}}, supports_credentials=True)

# Simple credentials for optional API login
USERNAME = os.environ.get('DOCUCHAT_USER', 'admin')
PASSWORD = os.environ.get('DOCUCHAT_PASS', 'password')

UPLOAD_FOLDER = 'static/uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Load your model and tokenizer here
MODEL_NAME = "bert-base-multilingual-cased"
MODEL_DIR = os.path.join("models", MODEL_NAME.replace("/", "_"))

if not os.path.exists(MODEL_DIR):
    raise ValueError(f"Model directory {MODEL_DIR} does not exist. Please run the download_models.py script to download the model.")

tokenizer = AutoTokenizer.from_pretrained(MODEL_DIR)
model = AutoModel.from_pretrained(MODEL_DIR)

def extract_text_from_docx(file_path):
    doc = DocxDocument(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_xlsx(file_path):
    workbook = load_workbook(filename=file_path, data_only=True)
    sheet = workbook.active
    return "\n".join(["\t".join([str(cell.value) for cell in row]) for row in sheet.iter_rows()])

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def handle_question(question, project_name):
    # Ensure the project-specific index and doc_list are loaded
    index, doc_list = load_project_index(project_name)

    # Tokenize the question and get its embeddings
    inputs = tokenizer(question, return_tensors='pt', truncation=True, padding=True, max_length=512)
    question_embedding = model(**inputs).last_hidden_state.mean(dim=1).detach().numpy()

    # Perform a search query
    D, I = index.search(np.array(question_embedding).astype('float32'), k=1)

    div_open = "<div>"
    div_close = "</div>"
    small_tag = "<small>Based on the embeddings, the best match for your question is:</small>"
    linebreak = "<br><br>"

    # Find the chunk with the highest similarity score
    if I[0][0] != -1:
        best_match = doc_list[I[0][0]]
        response = (
            "{div_open}{small_tag}{linebreak}"
            "{best_match}"
            "{div_close}"
        ).format(div_open=div_open, small_tag=small_tag, linebreak=linebreak, best_match=best_match.replace('\n', '<br>'), div_close=div_close)
    else:
        response = (
            "{div_open}{small_tag}{linebreak}"
            "No relevant information found."
            "{div_close}"
        ).format(div_open=div_open, small_tag=small_tag, linebreak=linebreak, div_close=div_close)

    return response

def load_project_index(project_name):
    project_dir = os.path.join("project_embeddings", project_name)
    if not os.path.exists(project_dir):
        raise ValueError(f"Project directory {project_dir} does not exist.")

    # Load FAISS index
    index_path = os.path.join(project_dir, "index.faiss")
    if not os.path.exists(index_path):
        raise ValueError(f"Index file {index_path} does not exist.")

    index = faiss.read_index(index_path)

    # Load document list
    doc_list_path = os.path.join(project_dir, "doc_list.json")
    if not os.path.exists(doc_list_path):
        raise ValueError(f"Document list file {doc_list_path} does not exist.")

    with open(doc_list_path, 'r') as f:
        doc_list = json.load(f)

    return index, doc_list

class TextDataset(torch.utils.data.Dataset):
    def __init__(self, texts, tokenizer, max_length):
        self.texts = texts
        self.tokenizer = tokenizer
        self.max_length = max_length

    def __len__(self):
        return len(self.texts)

    def __getitem__(self, idx):
        text = self.texts[idx]
        inputs = self.tokenizer(text, truncation=True, padding='max_length', max_length=self.max_length, return_tensors="pt")
        input_ids = inputs['input_ids'].squeeze()
        attention_mask = inputs['attention_mask'].squeeze()
        return {
            'input_ids': input_ids,
            'attention_mask': attention_mask,
            'labels': torch.tensor(1)  # Dummy label for compatibility
        }

@app.route('/upload', methods=['POST'])
def upload_document():
    if 'file' not in request.files:
        return jsonify({"status": "No file part"})

    file = request.files['file']
    if file.filename == '':
        return jsonify({"status": "No selected file"})

    project_name = request.form.get('project')
    if not project_name:
        return jsonify({"status": "No project name provided"})
    print(f"Uploading document to project: {project_name}")
    project_folder = os.path.join(app.config['UPLOAD_FOLDER'], project_name)
    if not os.path.exists(project_folder):
        os.makedirs(project_folder)

    file_path = os.path.join(project_folder, file.filename)
    file.save(file_path)
    print(f"File saved to {file_path}")

    # Extract text from the document and generate embeddings
    try:
        text = extract_text(file_path)

        print(f"Extracted text: {text[:100]}")

        chunk_embeddings, text_chunks = generate_embeddings_for_text(text)

        print(f"Embeddings generated for project: {project_name}")

        project_embedding_dir = os.path.join("project_embeddings", project_name)
        if not os.path.exists(project_embedding_dir):
            os.makedirs(project_embedding_dir)

        index = faiss.IndexFlatL2(768)
        index.add(np.array(chunk_embeddings).astype('float32'))
        faiss.write_index(index, os.path.join(project_embedding_dir, "index.faiss"))

        with open(os.path.join(project_embedding_dir, "doc_list.json"), 'w') as f:
            json.dump(text_chunks, f)

        return jsonify({"status": "success", "project": project_name})
    except Exception as e:
        print(f"Error during embedding generation: {e}")
        return jsonify({"status": "error", "message": str(e)})

def extract_text(file_path):
    if file_path.endswith('.pdf'):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith('.docx'):
        return extract_text_from_docx(file_path)
    elif file_path.endswith('.xlsx'):
        return extract_text_from_xlsx(file_path)
    elif file_path.endswith('.pptx') or file_path.endswith('.ppt'):
        return extract_text_from_pptx(file_path)
    elif file_path.endswith('.txt'):
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    else:
        raise ValueError("Unsupported file type")

def generate_embeddings_for_text(text):
    chunk_size = 512  # Set chunk size to 512 tokens
    tokens = tokenizer(text, return_tensors='pt', truncation=True, padding=True, max_length=chunk_size)
    chunk_embeddings = []
    text_chunks = []

    # Split text into chunks and process each chunk
    for i in range(0, len(tokens['input_ids'][0]), chunk_size):
        chunk_text = text[i:i + chunk_size]
        inputs = {k: v[:, i:i + chunk_size] for k, v in tokens.items()}
        outputs = model(**inputs).last_hidden_state.mean(dim=1).detach().numpy()
        chunk_embeddings.append(outputs[0])
        text_chunks.append(chunk_text)

    return chunk_embeddings, text_chunks

# Optional API login endpoint
@app.route('/login', methods=['POST'])
def api_login():
    data = request.json
    if not data:
        return jsonify({'status': 'failure'}), 400
    username = data.get('username')
    password = data.get('password')
    if username == USERNAME and password == PASSWORD:
        return jsonify({'status': 'success'})
    return jsonify({'status': 'failure'}), 401

@app.route('/projects/<project_name>/files', methods=['DELETE'])
def delete_file_from_project(project_name):
    data = request.json
    file_name = data.get('file')
    project_folder = os.path.join(app.config['UPLOAD_FOLDER'], project_name)
    file_path = os.path.join(project_folder, file_name)
    if os.path.exists(file_path):
        os.remove(file_path)
        return jsonify({"status": "success", "file": file_name})
    else:
        return jsonify({"status": "File not found"})

@app.route('/ask', methods=['POST'])
def ask_question():
    data = request.json
    question = data.get('question')
    project_name = data.get('project')
    print(f"Received question for project: {project_name}")

    session['history'] = session.get('history', []) + [question]
    response = handle_question(question, project_name)

    print(f"Response: {response}")
    return jsonify({"response": response})

@app.route('/projects', methods=['GET'])
def list_projects():
    projects = []
    for f in os.scandir(app.config['UPLOAD_FOLDER']):
        if f.is_dir():
            project_files = [file.name for file in os.scandir(f.path) if file.is_file()]
            projects.append({"project": f.name, "files": project_files})
    return jsonify({"projects": projects})

@app.route('/projects', methods=['POST'])
def create_project():
    data = request.json
    project_name = data.get('project')
    project_folder = os.path.join(app.config['UPLOAD_FOLDER'], project_name)
    if not os.path.exists(project_folder):
        os.makedirs(project_folder)
        return jsonify({"status": "success", "project": project_name})
    else:
        return jsonify({"status": "Project already exists"})

@app.route('/projects', methods=['DELETE'])
def delete_project():
    data = request.json
    project_name = data.get('project')
    project_folder = os.path.join(app.config['UPLOAD_FOLDER'], project_name)
    if os.path.exists(project_folder):
        for file in os.scandir(project_folder):
            os.remove(file.path)
        os.rmdir(project_folder)
        project_embedding_dir = os.path.join("project_embeddings", project_name)
        if os.path.exists(project_embedding_dir):
            for file in os.scandir(project_embedding_dir):
                os.remove(file.path)
            os.rmdir(project_embedding_dir)
        return jsonify({"status": "success", "project": project_name})
    else:
        return jsonify({"status": "Project not found"})

@app.route('/projects/<project_name>/generate_embeddings', methods=['POST'])
def generate_embeddings(project_name):
    project_folder = os.path.join(app.config['UPLOAD_FOLDER'], project_name)
    if not os.path.exists(project_folder):
        return jsonify({"status": "Project not found"})

    try:
        text_chunks = []
        chunk_embeddings = []
        for file in os.listdir(project_folder):
            file_path = os.path.join(project_folder, file)
            text = extract_text(file_path)

            print(f"Extracted text: {text[:100]}")

            new_chunk_embeddings, new_text_chunks = generate_embeddings_for_text(text)
            chunk_embeddings.extend(new_chunk_embeddings)
            text_chunks.extend(new_text_chunks)

        print(f"Embeddings generated for project: {project_name}")

        project_embedding_dir = os.path.join("project_embeddings", project_name)
        if not os.path.exists(project_embedding_dir):
            os.makedirs(project_embedding_dir)

        index = faiss.IndexFlatL2(768)
        index.add(np.array(chunk_embeddings).astype('float32'))
        faiss.write_index(index, os.path.join(project_embedding_dir, "index.faiss"))

        with open(os.path.join(project_embedding_dir, "doc_list.json"), 'w') as f:
            json.dump(text_chunks, f)

        return jsonify({"status": "success", "project": project_name})
    except Exception as e:
        print(f"Error during embedding generation: {e}")
        return jsonify({"status": "error", "message": str(e)})

@app.route('/fine_tune', methods=['POST'])
def fine_tune():
    project_name = request.json.get('project_name', 'Morisco')  # Adjust this as necessary
    project_embedding_dir = os.path.join("project_embeddings", project_name)
    if not os.path.exists(project_embedding_dir):
        return jsonify({"status": "Project not found"})

    try:
        with open(os.path.join(project_embedding_dir, "doc_list.json"), 'r') as f:
            texts = json.load(f)

        labels = [0] * len(texts)  # Dummy labels, adjust if you have actual labels

        fine_tuner = ModelFineTuner()
        fine_tuner.fine_tune(texts, labels)

        return jsonify({"status": "success"})
    except Exception as e:
        return jsonify({"status": "error", "message": str(e)})

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=8080, debug=True)
